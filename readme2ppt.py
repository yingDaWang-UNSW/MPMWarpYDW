"""
Create MPMWarpYDW PowerPoint presentation with proper math equations.
Run: python create_pptx_with_math.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=''):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(102, 102, 102)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.7))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(point, tuple):
            p.text = point[0]
            p.level = point[1]
        else:
            p.text = point
            p.level = 0
        
        p.font.size = Pt(18) if p.level == 0 else Pt(16)
        p.space_after = Pt(8)
    
    return slide

def add_math_slide(prs, title, items):
    """Add slide with math equations (using Cambria Math font)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    y_pos = 1.2
    for desc, eq in items:
        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.bold = True
        y_pos += 0.45
        
        # Equation in Cambria Math
        eq_box = slide.shapes.add_textbox(Inches(1.0), Inches(y_pos), Inches(11.333), Inches(0.5))
        tf = eq_box.text_frame
        p = tf.paragraphs[0]
        p.text = eq
        p.font.size = Pt(22)
        p.font.name = 'Cambria Math'
        p.font.color.rgb = RGBColor(0, 0, 139)
        y_pos += 0.65
    
    return slide

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, 
    'MPMWarpYDW',
    'GPU-Accelerated Coupled MPM-XPBD Framework\nfor Geomaterial Failure and Fragmentation')

# Slide 2: Overview
add_content_slide(prs, 'Framework Overview', [
    'Hybrid continuum-discrete simulation for progressive material failure:',
    ('MPM Phase: Elasto-plastic deformation with damage accumulation', 1),
    ('Phase Transition: Automatic MPM → XPBD when damage ≥ 1.0', 1),
    ('XPBD Phase: Discrete particle dynamics with contact resolution', 1),
    '',
    'Target Applications:',
    ('Rock avalanches and landslides', 1),
    ('Cave propagation and block caving', 1),
    ('Dynamic fragmentation under impact', 1),
    ('Granular flow with cohesion', 1),
])

# Slide 3: MPM Features
add_content_slide(prs, 'Material Point Method (MPM) Features', [
    'Constitutive Framework:',
    ('Multiplicative elastoplasticity with logarithmic (Hencky) strain', 1),
    ('Von Mises (J₂) and Drucker-Prager plasticity models', 1),
    ('Return mapping in principal strain space', 1),
    '',
    'Damage & Phase Transition:',
    ('Isotropic scalar damage from accumulated plastic strain', 1),
    ('Automatic transition to XPBD at critical damage', 1),
    ('Energy release: elastic strain energy → kinetic energy', 1),
    '',
    'Numerical Schemes:',
    ('APIC transfer (angular momentum conservation)', 1),
    ('Optional RPIC for controllable damping', 1),
    ('Kelvin-Voigt viscoelasticity for stabilization', 1)
])

# Slide 4: Kinematic Setup
add_math_slide(prs, 'Kinematic Framework', [
    ('Deformation gradient update:', 'F_trial = (I + ∇v·Δt) · Fₙ'),
    ('SVD decomposition:', 'F = U · Σ · Vᵀ'),
    ('Logarithmic (Hencky) strain:', 'ε = log(Σ) = [log(σ₁), log(σ₂), log(σ₃)]'),
    ('Volumetric-deviatoric split:', 'ε_vol = tr(ε)/3,    ε_dev = ε − ε_vol·I'),
])

# Slide 5: Stress Computation
add_math_slide(prs, 'Stress Computation', [
    ('Kirchhoff stress:', 'τ = 2μ·ε + λ·tr(ε)·I'),
    ('Shear modulus:', 'μ = E / [2(1 + ν)]'),
    ('First Lamé parameter:', 'λ = E·ν / [(1 + ν)(1 − 2ν)]'),
    ('Cauchy stress:', 'σ = τ / det(F)'),
])

# Slide 6: Von Mises
add_math_slide(prs, 'Von Mises Plasticity (J₂)', [
    ('Yield criterion:', 'f = ‖τ_dev‖ − √(2/3)·(1−D)·σ_y ≤ 0'),
    ('Associated flow rule:', 'Δε_plastic = Δγ · (τ_dev / ‖τ_dev‖)'),
    ('Plastic multiplier:', 'Δγ = ‖ε_dev‖ − σ_y,eff / (2μ)'),
])

# Slide 7: Drucker-Prager
add_math_slide(prs, 'Drucker-Prager Plasticity', [
    ('Yield criterion:', 'f = ‖τ_dev‖ − √(2/3)·(1−D)·(σ_y − α·p) ≤ 0'),
    ('Mean stress:', 'p = tr(τ)/3    (compression negative)'),
    ('Non-associated flow:', 'Δε_plastic = Δγ · [n_dev + β·I/3]'),
    ('Dilatancy parameter:', 'β = 0.3·α'),
])

# Slide 8: Damage
add_math_slide(prs, 'Damage Evolution & Phase Transition', [
    ('Damage accumulation:', 'D = min(1.0, Σ(Δε_plastic) / ε_critical)'),
    ('Strength degradation:', 'σ_y,eff = (1 − D) · σ_y'),
    ('Strain energy density:', 'u = ½ · τ : ε'),
    ('Velocity at transition:', 'v_release = √(2·η·u/ρ)'),
])

# Slide 9: XPBD Features
add_content_slide(prs, 'Extended Position Based Dynamics (XPBD)', [
    'Contact Resolution:',
    ('Spatial hash grid for O(n) neighbor search', 1),
    ('Iterative Gauss-Seidel constraint projection', 1),
    ('SOR relaxation for faster convergence', 1),
    '',
    'Friction Model:',
    ('Coulomb friction with static/dynamic transition', 1),
    ('Velocity-based threshold for static friction', 1),
    '',
    'Additional Features:',
    ('Particle cohesion (cementation)', 1),
    ('Particle sleeping (computational efficiency)', 1),
    ('Particle swelling (fragmentation bulking)', 1)
])

# Slide 10: XPBD Equations
add_math_slide(prs, 'XPBD Contact Mechanics', [
    ('Contact constraint:', 'C = ‖xᵢ − xⱼ‖ − (rᵢ + rⱼ) ≥ 0'),
    ('Normal correction:', 'Δx_n = C · n · wᵢ / (wᵢ + wⱼ)'),
    ('Contact normal:', 'n = (xᵢ − xⱼ) / ‖xᵢ − xⱼ‖'),
    ('Friction (Coulomb):', 'Δx_f = clamp(Δv_t·Δt, −μ·|Δx_n|, μ·|Δx_n|)'),
])

# Slide 11: Coupling
add_content_slide(prs, 'MPM-XPBD Coupling Strategy', [
    'Contact-Based Momentum Transfer:',
    ('Direct particle-particle contact forces', 1),
    ('Same XPBD constraint projection for MPM-XPBD contacts', 1),
    ('xpbd_mpm_coupling_strength scales impulse magnitude', 1),
    '',
    'Material Label System:',
    ('Label 0: MPM in contact with XPBD (cannot transition)', 1),
    ('Label 1: MPM not in contact (can transition if D ≥ 1)', 1),
    ('Label 2: XPBD discrete particles', 1),
    '',
    'Transition Locking (prevents cascading failure):',
    ('Reset: all label 0 → label 1 before XPBD step', 1),
    ('Detect: if MPM touches XPBD, set label → 0', 1),
    ('Transition: only label 1 can become XPBD', 1)
])

# Slide 12: Geostatic
add_math_slide(prs, 'Geostatic Initialization (K₀)', [
    ('Vertical stress:', 'σ_v = −ρ·g·(z_top − z)'),
    ('Horizontal stress:', 'σ_h = K₀·σ_v'),
    ('Initialize F (NOT stress):', 'F_init = diag(exp(ε₁), exp(ε₂), exp(ε₃))'),
])

# Slide 13: Viscoelasticity
add_math_slide(prs, 'Viscoelasticity (Kelvin-Voigt)', [
    ('Viscous stress:', 'τ_visc = 2·η_shear·ε̇_dev + η_bulk·tr(ε̇)·I'),
    ('Strain rate tensor:', 'ε̇ = ½·(∇v + ∇vᵀ)'),
    ('Typical values:', 'η = 10³ − 10⁶ Pa·s'),
])

# Slide 14: CFL
add_math_slide(prs, 'Numerical Stability', [
    ('CFL condition:', 'dt < α · dx / c_p    (α ≈ 0.3−0.5)'),
    ('P-wave speed:', 'c_p = √[(K + 4μ/3) / ρ]'),
    ('Cell-crossing limit:', 'max(‖Δx‖) < 0.5 · dx'),
])

# Slide 15: Simulation Loop
add_content_slide(prs, 'Simulation Loop Structure', [
    'Big Step Loop (outer):',
    ('Phase 1: Combined MPM+XPBD simulation', 1),
    ('Phase 2: XPBD-only settling (optional)', 1),
    ('Checkpoint save at big step boundary', 1),
    '',
    'Phase 1 Inner Loop:',
    ('MPM: stress → P2G → grid ops → G2P (every dt)', 1),
    ('XPBD: hash grid → contacts → integrate (every dtxpbd)', 1),
    ('Early termination: damage stall or convergence', 1),
    '',
    'Phase 2 (if xpbdOnlyDuration > 0):',
    ('XPBD steps with frozen MPM', 1),
    ('Sleep-based early termination', 1)
])

# Slide 16: Parameters
add_content_slide(prs, 'Key Configuration Parameters', [
    'Time Stepping:',
    ('dt: MPM timestep (0 = auto CFL)', 1),
    ('dtxpbd: XPBD timestep (10-100× dt typical)', 1),
    ('bigStepDuration, xpbdOnlyDuration: phase durations', 1),
    '',
    'Material:',
    ('E, nu: elastic moduli', 1),
    ('ys, alpha: yield stress, pressure sensitivity', 1),
    ('strainCriteria: plastic strain for full damage', 1),
    ('eta_shear, eta_bulk: viscosity coefficients', 1),
    '',
    'Coupling:',
    ('xpbd_mpm_coupling_strength: impulse scaling', 1),
    ('mpm_contact_transition_lock: prevent cascading', 1)
])

# Slide 17: Output
add_content_slide(prs, 'Output Files', [
    'Particle VTP Files (ParaView time series):',
    ('positions, radius, velocity', 1),
    ('damage, mean_stress, von_mises', 1),
    ('stress_tensor (6 components: xx, yy, zz, xy, xz, yz)', 1),
    ('active_label, material_label', 1),
    '',
    'Checkpoint Data (at big step ends):',
    ('Full deformation gradient F, APIC matrix C', 1),
    ('All material properties per particle', 1),
    ('Enables restart: --restart path/to/checkpoint.vtp', 1),
])

# Slide 18: Performance
add_content_slide(prs, 'Performance Considerations', [
    'GPU Scaling:',
    ('Particle operations: O(N) linear', 1),
    ('Grid operations: O(G³) with grid volume', 1),
    ('Contacts: O(N) with spatial hashing', 1),
    '',
    'Memory (approximate):',
    ('~500 bytes/particle, ~100 bytes/grid cell', 1),
    ('1M particles + 256³ grid ≈ 2.5 GB VRAM', 1),
    '',
    'Tuning Tips:',
    ('Stiff materials: small dt, high xpbd_iterations', 1),
    ('Phase transitions: use viscosity, eff < 1.0', 1),
    ('Large domains: grid_particle_spacing_scale = 2.0', 1)
])

# Slide 19: References
add_content_slide(prs, 'Key References', [
    'MPM:',
    ('Sulsky et al. (1994) - Original MPM formulation', 1),
    ('Stomakhin et al. (2013) - Snow simulation', 1),
    '',
    'Transfer Schemes:',
    ('Jiang et al. (2015) - APIC', 1),
    ('Jiang et al. (2017) - SIGGRAPH MPM Course', 1),
    '',
    'XPBD:',
    ('Macklin et al. (2016) - XPBD formulation', 1),
    '',
    'Constitutive:',
    ('de Souza Neto et al. (2011) - Computational Plasticity', 1),
    ('Simo (1992) - Multiplicative plasticity', 1)
])

# Slide 20: Summary
add_title_slide(prs,
    'Summary',
    'GPU-accelerated MPM-XPBD coupling\nMultiplicative elastoplasticity with damage\nAutomatic continuum → discrete transition\nApplications: geomaterial failure, fragmentation, granular flow')

# Save
prs.save('MPMWarpYDW_Presentation.pptx')
print('Saved: MPMWarpYDW_Presentation.pptx')